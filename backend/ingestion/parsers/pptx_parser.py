"""PPTX Parser — extracts text from PowerPoint presentations."""

import io
from pptx import Presentation
from ingestion import TextChunk
from utils.logger import logger


def parse_pptx(content: bytes, filename: str) -> list[TextChunk]:
    """Parse a PPTX file, extracting text from each slide."""
    chunks = []
    try:
        prs = Presentation(io.BytesIO(content))
        total_slides = len(prs.slides)
        logger.info(f"Parsing PPTX: {filename} ({total_slides} slides)")

        for i, slide in enumerate(prs.slides):
            slide_texts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                # Extract table content
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))

            # Include speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Speaker Notes] {notes}")

            if slide_texts:
                chunks.append(TextChunk(
                    text="\n".join(slide_texts),
                    page_number=i + 1,
                    section=f"Slide {i + 1}",
                    metadata={"source": filename, "total_slides": total_slides},
                ))

        logger.info(f"PPTX parsed: {len(chunks)} slides with content")
    except Exception as e:
        logger.error(f"PPTX parsing failed for {filename}: {e}")
        raise

    return chunks
